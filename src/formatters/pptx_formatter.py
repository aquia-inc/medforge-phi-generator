"""
PPTX (PowerPoint) formatter for PHI documents
Creates presentations with case studies and de-identified educational content
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
import os


class PPTXFormatter:
    """Creates PowerPoint presentations with PHI content"""

    def __init__(self, output_dir='output'):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def create_case_study_presentation(self, patient, provider, facility, filename):
        """Create case study presentation (PHI Positive)"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Clinical Case Study"
        subtitle.text = f"{facility['name']}\n{datetime.now().strftime('%B %Y')}"

        # Slide 2: Patient Information
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        title = slide.shapes.title
        title.text = "Patient Information"

        # Add text box with patient details
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        patient_info = text_frame.add_paragraph()
        patient_info.text = f"Patient: {patient['first_name']} {patient['last_name']}\n"
        patient_info.text += f"MRN: {patient['mrn']}\n"
        patient_info.text += f"DOB: {patient['dob'].strftime('%m/%d/%Y')} (Age: {patient['age']})\n"
        patient_info.text += f"Contact: {patient['phone']} / {patient['email']}\n"
        patient_info.text += f"Address: {patient['address']['street']}, {patient['address']['city']}, {patient['address']['state']}"
        patient_info.font.size = Pt(18)

        # Slide 3: Medical History
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Medical History"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        history = text_frame.add_paragraph()
        history.text = "Diagnoses:\n"
        history.font.bold = True
        history.font.size = Pt(16)

        for diag in patient['diagnoses']:
            p = text_frame.add_paragraph()
            p.text = f"• {diag['name']} (ICD-10: {diag['icd10']})"
            p.level = 1
            p.font.size = Pt(14)

        meds = text_frame.add_paragraph()
        meds.text = "\nCurrent Medications:"
        meds.font.bold = True
        meds.font.size = Pt(16)

        for med in patient['medications'][:5]:
            p = text_frame.add_paragraph()
            p.text = f"• {med}"
            p.level = 1
            p.font.size = Pt(14)

        # Slide 4: Vital Signs
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Vital Signs"

        vitals = patient['vital_signs']
        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(7)
        height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        vital_text = (
            f"Blood Pressure: {vitals['blood_pressure']} mmHg\n"
            f"Heart Rate: {vitals['heart_rate']} bpm\n"
            f"Temperature: {vitals['temperature']}°F\n"
            f"Respiratory Rate: {vitals['respiratory_rate']}\n"
            f"O2 Saturation: {vitals['oxygen_saturation']}%\n"
            f"Weight: {vitals['weight']} lbs\n"
            f"Height: {vitals['height']} inches"
        )

        p = text_frame.add_paragraph()
        p.text = vital_text
        p.font.size = Pt(20)

        # Slide 5: Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Clinical Summary"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        summary = text_frame.add_paragraph()
        summary.text = f"Patient under care of {provider['first_name']} {provider['last_name']}, {provider['title']}\n"
        summary.text += f"Specialty: {provider['specialty']}\n\n"
        summary.text += f"Primary diagnosis: {patient['diagnoses'][0]['name']}\n"
        summary.text += "Management plan includes medication therapy and regular follow-up."
        summary.font.size = Pt(16)

        # Footer with confidentiality notice
        footer = text_frame.add_paragraph()
        footer.text = "\n\nCONFIDENTIAL: Contains Protected Health Information"
        footer.font.size = Pt(12)
        footer.font.italic = True
        footer.font.color.rgb = RGBColor(255, 0, 0)

        # Save
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        return filepath

    def create_educational_presentation(self, facility, filename):
        """Create educational presentation (PHI Negative - No Patient Data)"""
        prs = Presentation()

        # Slide 1: Title
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Diabetes Management Guidelines"
        subtitle.text = f"{facility['name']}\nClinical Education Series"

        # Slide 2: Overview
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Type 2 Diabetes Overview"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        overview = text_frame.add_paragraph()
        overview.text = "Key Points:"
        overview.font.bold = True
        overview.font.size = Pt(18)

        points = [
            "Prevalence: Affects approximately 34 million Americans",
            "Risk Factors: Obesity, sedentary lifestyle, family history",
            "Complications: Cardiovascular disease, neuropathy, nephropathy",
            "Management: Lifestyle modification + medication as needed"
        ]

        for point in points:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 1
            p.font.size = Pt(16)

        # Slide 3: Treatment Guidelines
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Treatment Approach"

        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        treatment = text_frame.add_paragraph()
        treatment.text = "Step-wise Approach:"
        treatment.font.bold = True
        treatment.font.size = Pt(18)

        steps = [
            "Step 1: Lifestyle modifications (diet, exercise)",
            "Step 2: Metformin as first-line medication",
            "Step 3: Add second agent if A1C > 7%",
            "Step 4: Consider insulin for A1C > 9%",
            "Regular monitoring: A1C every 3 months"
        ]

        for step in steps:
            p = text_frame.add_paragraph()
            p.text = step
            p.level = 1
            p.font.size = Pt(16)

        # Slide 4: Monitoring
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Laboratory Monitoring"

        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(7)
        height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        monitoring = (
            "Hemoglobin A1C: Every 3 months\n"
            "Lipid Panel: Annually\n"
            "Kidney Function: Annually\n"
            "Retinal Exam: Annually\n"
            "Foot Exam: Each visit"
        )

        p = text_frame.add_paragraph()
        p.text = monitoring
        p.font.size = Pt(18)

        # Slide 5: Summary
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Summary"

        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        summary = text_frame.add_paragraph()
        summary.text = "Early detection and comprehensive management improve outcomes.\n\n"
        summary.text += "Multidisciplinary approach recommended:\n"
        summary.text += "Primary care, endocrinology, nutrition, ophthalmology"
        summary.font.size = Pt(18)

        footer = text_frame.add_paragraph()
        footer.text = f"\n\n{facility['name']} Clinical Guidelines"
        footer.font.size = Pt(12)
        footer.font.italic = True

        # Save
        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        return filepath
