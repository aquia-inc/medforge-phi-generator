"""
CUI PPTX formatter.

Creates PowerPoint presentations with CUI markings and content,
mirroring the PHI PPTXFormatter capabilities for CUI format parity.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import os
from typing import Any, Dict, Optional

from templates.components import ComponentConfiguration, get_docx_font_name

# LLM narrative fields that get dedicated slides when present
_NARRATIVE_SLIDES = [
    ('executive_summary', 'Executive Summary'),
    ('body_content', 'Details'),
    ('analysis', 'Analysis'),
    ('risk_assessment', 'Risk Assessment'),
    ('recommendations', 'Recommendations'),
    ('justification', 'Justification'),
]


class CUIPPTXFormatter:
    """Creates PowerPoint presentations with CUI content."""

    def __init__(self, output_dir: str = 'output'):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def create_cui_presentation(self, doc_data: Dict[str, Any], filename: str,
                                component_config: Optional[ComponentConfiguration] = None) -> str:
        """
        Create a PPTX presentation from CUI data with optional visual variation.

        Slides: Title (classification banner + title), Metadata (agency, date,
        authority), Content (2-3 slides based on doc_data sections),
        LLM narrative slides (if present), Footer (distribution + classification).

        Args:
            doc_data: Dictionary containing CUI document data
            filename: Output filename
            component_config: Optional component configuration for visual variety

        Returns:
            Path to created file
        """
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Extract style settings
        style_cfg = component_config.style.get_config() if component_config else None
        header_cfg = component_config.header.get_config() if component_config else None

        # Resolve font name (PPTX uses Windows names like DOCX)
        font_name = get_docx_font_name(style_cfg["font_family"]) if style_cfg else None
        title_size = style_cfg.get("font_size_title", 16) if style_cfg else 16
        body_size = style_cfg.get("font_size_body", 14) if style_cfg else 14
        primary_color = None
        if style_cfg and style_cfg.get("use_colors"):
            try:
                c = style_cfg["color_primary"].lstrip("#")
                primary_color = RGBColor(int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16))
            except (ValueError, IndexError, KeyError):
                pass

        # Bundle style for passing to helper methods
        pptx_style = {
            'font_name': font_name,
            'title_size': title_size,
            'body_size': body_size,
            'primary_color': primary_color,
        }

        classification = doc_data.get('classification', '')
        has_cui = doc_data.get('has_cui', False)
        title_text = doc_data.get('title', 'Document')

        # Slide 1: Title slide with classification banner
        self._add_title_slide(prs, title_text, classification, has_cui, doc_data, pptx_style)

        # Slide 2: Metadata
        self._add_metadata_slide(prs, doc_data, pptx_style)

        # Slides 3+: Content slides
        self._add_content_slides(prs, doc_data, pptx_style)

        # LLM narrative slides (when LLM-enhanced fields are present)
        self._add_narrative_slides(prs, doc_data, pptx_style)

        # Final slide: Distribution / classification footer
        if has_cui:
            self._add_footer_slide(prs, doc_data, classification)

        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)
        return filepath

    def _add_title_slide(self, prs, title_text, classification, has_cui, doc_data,
                          pptx_style=None):
        """Add title slide with optional classification banner."""
        s = pptx_style or {}
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = title_text
        # Apply style to title
        for run in title.text_frame.paragraphs[0].runs:
            if s.get('font_name'):
                run.font.name = s['font_name']
            if s.get('primary_color'):
                run.font.color.rgb = s['primary_color']

        agency = doc_data.get('agency', '')
        date_str = doc_data.get('document_date', datetime.now().strftime('%B %Y'))
        subtitle.text = f"{agency}\n{date_str}"

        if has_cui and classification:
            txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = classification
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(139, 0, 0)
            p.alignment = PP_ALIGN.CENTER

    def _add_metadata_slide(self, prs, doc_data, pptx_style=None):
        """Add metadata slide with agency, date, authority."""
        s = pptx_style or {}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Document Information"

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame

        fields = [
            ('Agency', doc_data.get('agency', '')),
            ('Date', doc_data.get('document_date', '')),
            ('Authority', doc_data.get('authority', '')),
            ('Category', doc_data.get('category', '').replace('_', ' ').title()),
            ('Document Type', doc_data.get('document_type', '').replace('_', ' ').title()),
        ]

        label_size = Pt(s.get('title_size', 16))
        body_size = Pt(s.get('body_size', 16))

        for label, value in fields:
            if value:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{label}: "
                run.font.bold = True
                run.font.size = label_size
                if s.get('font_name'):
                    run.font.name = s['font_name']
                run2 = p.add_run()
                run2.text = str(value)
                run2.font.size = body_size
                if s.get('font_name'):
                    run2.font.name = s['font_name']

    def _add_content_slides(self, prs, doc_data, pptx_style=None):
        """Add content slides based on document type."""
        doc_type = doc_data.get('document_type', '')

        if doc_type == 'vulnerability_alert':
            self._add_vulnerability_slides(prs, doc_data, pptx_style)
        elif doc_type in ('budget_memo', 'eft_authorization', 'retirement_estimate'):
            self._add_financial_slides(prs, doc_data, pptx_style)
        elif doc_type in ('investigation_summary', 'criminal_history_check'):
            self._add_law_enforcement_slides(prs, doc_data, pptx_style)
        else:
            self._add_generic_content_slides(prs, doc_data, pptx_style)

    def _add_vulnerability_slides(self, prs, doc_data, pptx_style=None):
        """Add vulnerability alert content slides."""
        s = pptx_style or {}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Vulnerability Details"

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame

        fields = [
            ('Severity', doc_data.get('severity', '')),
            ('CVSS Score', doc_data.get('cvss_score', '')),
            ('CVE', doc_data.get('cve_id', '')),
            ('Affected System', doc_data.get('affected_system', '')),
        ]

        for label, value in fields:
            if value:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{label}: {value}"
                run.font.size = Pt(s.get('body_size', 16))
                if s.get('font_name'):
                    run.font.name = s['font_name']

        desc = doc_data.get('description', '')
        if desc:
            self._add_text_slide(prs, "Description & Remediation", desc, s)
            rem = doc_data.get('remediation', {})
            if isinstance(rem, dict) and rem.get('action'):
                action_text = f"Remediation: {rem['action']}"
                if rem.get('deadline'):
                    action_text += f"\nDeadline: {rem['deadline']}"
                self._add_text_slide(prs, "Remediation Plan", action_text, s)

    def _add_financial_slides(self, prs, doc_data, pptx_style=None):
        """Add financial document content slides."""
        s = pptx_style or {}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = doc_data.get('title', 'Financial Details')

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True

        skip_fields = {'document_id', 'document_type', 'category', 'subcategory',
                       'has_cui', 'classification', 'generated_date', 'title',
                       'confidentiality_notice', 'document_date', 'agency', 'authority',
                       'executive_summary', 'body_content', 'recommendations',
                       'risk_assessment', 'analysis', 'justification'}

        for key, value in doc_data.items():
            if key in skip_fields:
                continue
            if isinstance(value, (str, int, float)) and value:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{key.replace('_', ' ').title()}: {value}"
                run.font.size = Pt(s.get('body_size', 14))
                if s.get('font_name'):
                    run.font.name = s['font_name']

    def _add_law_enforcement_slides(self, prs, doc_data, pptx_style=None):
        """Add law enforcement content slides."""
        self._add_financial_slides(prs, doc_data, pptx_style)

    def _add_generic_content_slides(self, prs, doc_data, pptx_style=None):
        """Add generic content slides for any document type."""
        s = pptx_style or {}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Details"

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True

        skip_fields = {'document_id', 'document_type', 'category', 'subcategory',
                       'has_cui', 'classification', 'generated_date', 'title',
                       'confidentiality_notice', 'document_date', 'agency', 'authority',
                       'executive_summary', 'body_content', 'recommendations',
                       'risk_assessment', 'analysis', 'justification'}

        for key, value in doc_data.items():
            if key in skip_fields:
                continue
            if isinstance(value, (str, int, float)) and value:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = f"{key.replace('_', ' ').title()}: {value}"
                run.font.size = Pt(s.get('body_size', 14))
                if s.get('font_name'):
                    run.font.name = s['font_name']

        desc = doc_data.get('description', '')
        if desc and len(desc) > 200:
            self._add_text_slide(prs, "Additional Details", desc, s)

    def _add_text_slide(self, prs, title_text: str, body_text: str, pptx_style=None):
        """Add a slide with a title and multi-paragraph body text."""
        s = pptx_style or {}
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_text

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Split on double-newline for multi-paragraph LLM content
        paragraphs = [p.strip() for p in body_text.split('\n\n') if p.strip()]
        if not paragraphs:
            paragraphs = [body_text]

        for para_text in paragraphs:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = para_text
            run.font.size = Pt(s.get('body_size', 14))
            if s.get('font_name'):
                run.font.name = s['font_name']

    def _add_narrative_slides(self, prs, doc_data: Dict[str, Any], pptx_style=None):
        """Add dedicated slides for LLM-enriched narrative fields when present."""
        for field_key, slide_title in _NARRATIVE_SLIDES:
            value = doc_data.get(field_key, '')
            if isinstance(value, str) and len(value) > 30:
                self._add_text_slide(prs, slide_title, value, pptx_style)

    def _add_footer_slide(self, prs, doc_data, classification):
        """Add distribution/classification footer slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

        # Classification banner
        if classification:
            txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(0.5))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = classification
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(139, 0, 0)
            p.alignment = 1

        # Distribution notice
        notice = doc_data.get('confidentiality_notice', '')
        if notice:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.add_paragraph()
            p2.text = "DISTRIBUTION NOTICE"
            p2.font.size = Pt(18)
            p2.font.bold = True

            p3 = tf2.add_paragraph()
            p3.text = notice
            p3.font.size = Pt(12)
