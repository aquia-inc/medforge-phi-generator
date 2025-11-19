"""
PHI Document Validation System

Validates synthetic PHI documents across multiple formats:
- DOCX (Word documents)
- PDF (Portable Document Format)
- XLSX (Excel spreadsheets)
- EML (Email messages)
- PPTX (PowerPoint presentations)

Ensures PHI-positive documents contain required identifiers and
PHI-negative documents contain no patient identifiers.
"""

import re
import json
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Set
from dataclasses import dataclass, asdict
import zipfile

# Document format handlers
try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import extract_msg
except ImportError:
    extract_msg = None


@dataclass
class PHIElement:
    """Represents a detected PHI element"""
    element_type: str  # 'name', 'dob', 'mrn', 'ssn', 'address', 'phone', 'email'
    value: str
    location: str  # Where it was found (page, sheet, slide number, etc.)
    confidence: float = 1.0


@dataclass
class ValidationResult:
    """Results from validating a document"""
    filepath: str
    is_valid: bool
    file_format: str
    file_integrity_ok: bool
    expected_phi_type: str  # 'positive', 'negative', or 'unknown'
    phi_elements_found: List[PHIElement]
    validation_errors: List[str]
    validation_warnings: List[str]
    metadata: Dict
    validated_at: str

    def to_dict(self):
        """Convert to dictionary for JSON serialization"""
        result = asdict(self)
        result['phi_elements_found'] = [asdict(elem) for elem in self.phi_elements_found]
        return result


class PHIPatterns:
    """Regular expression patterns for detecting PHI elements"""

    # SSN pattern: ###-##-####
    SSN = re.compile(r'\b\d{3}-\d{2}-\d{4}\b')

    # MRN pattern: MRN followed by 6 digits
    MRN = re.compile(r'\bMRN\d{6}\b', re.IGNORECASE)

    # Phone pattern: (###) ###-####
    PHONE = re.compile(r'\(\d{3}\)\s*\d{3}-\d{4}')

    # Alternative phone patterns
    PHONE_ALT = re.compile(r'\b\d{3}[-.\s]?\d{3}[-.\s]?\d{4}\b')

    # DOB pattern: ##/##/####
    DOB = re.compile(r'\b\d{1,2}/\d{1,2}/\d{4}\b')

    # Email pattern
    EMAIL = re.compile(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b')

    # Address pattern (street number and name)
    ADDRESS = re.compile(r'\b\d+\s+[A-Z][a-z]+\s+(Street|St|Avenue|Ave|Road|Rd|Boulevard|Blvd|Lane|Ln|Drive|Dr|Court|Ct|Way|Place|Pl)\b', re.IGNORECASE)

    # Name pattern (capitalized first and last name)
    # Looking for "Patient:" or "Name:" followed by capitalized words
    NAME = re.compile(r'\b(?:Patient|Name):\s*([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)\b')

    # Insurance ID pattern
    INSURANCE_ID = re.compile(r'\b(?:MC|MD|BC|UH|AE)\d{8}\b')

    # ZIP code pattern
    ZIP_CODE = re.compile(r'\b\d{5}(?:-\d{4})?\b')


class PHIValidator:
    """Main validator class for PHI documents"""

    def __init__(self):
        self.patterns = PHIPatterns()
        self.supported_formats = {'.docx', '.pdf', '.xlsx', '.eml', '.pptx'}

    def validate_document(self, filepath: str, expected_phi_type: str = 'unknown') -> ValidationResult:
        """
        Validate a document for PHI content and integrity

        Args:
            filepath: Path to the document to validate
            expected_phi_type: 'positive', 'negative', or 'unknown'

        Returns:
            ValidationResult object with complete validation details
        """
        filepath = Path(filepath)
        errors = []
        warnings = []
        phi_elements = []
        metadata = {}

        # Check file exists
        if not filepath.exists():
            errors.append(f"File does not exist: {filepath}")
            return ValidationResult(
                filepath=str(filepath),
                is_valid=False,
                file_format='unknown',
                file_integrity_ok=False,
                expected_phi_type=expected_phi_type,
                phi_elements_found=[],
                validation_errors=errors,
                validation_warnings=warnings,
                metadata={},
                validated_at=datetime.now().isoformat()
            )

        # Get file format
        file_format = filepath.suffix.lower()
        if file_format not in self.supported_formats:
            errors.append(f"Unsupported file format: {file_format}")

        # Check file integrity
        integrity_ok = self.check_file_integrity(str(filepath))
        if not integrity_ok:
            errors.append("File integrity check failed - file may be corrupted")

        # Extract metadata
        try:
            metadata = self._extract_metadata(filepath)
        except Exception as e:
            warnings.append(f"Could not extract metadata: {str(e)}")

        # Extract PHI elements
        try:
            phi_elements = self.extract_phi_elements(str(filepath))
        except Exception as e:
            errors.append(f"Error extracting PHI elements: {str(e)}")

        # Validate based on expected type
        is_valid = True
        if expected_phi_type == 'positive':
            is_valid = self._validate_phi_positive(phi_elements, errors, warnings)
        elif expected_phi_type == 'negative':
            is_valid = self._validate_phi_negative(phi_elements, errors, warnings)

        return ValidationResult(
            filepath=str(filepath),
            is_valid=is_valid and integrity_ok and len(errors) == 0,
            file_format=file_format,
            file_integrity_ok=integrity_ok,
            expected_phi_type=expected_phi_type,
            phi_elements_found=phi_elements,
            validation_errors=errors,
            validation_warnings=warnings,
            metadata=metadata,
            validated_at=datetime.now().isoformat()
        )

    def check_phi_positive(self, filepath: str, expected_phi_elements: Set[str] = None) -> bool:
        """
        Check if a PHI-positive document contains required identifiers

        Args:
            filepath: Path to document
            expected_phi_elements: Set of required element types
                                 {'name', 'dob', 'mrn', 'ssn', 'address', 'phone'}

        Returns:
            True if all required elements are present
        """
        if expected_phi_elements is None:
            expected_phi_elements = {'name', 'dob', 'mrn', 'ssn', 'address', 'phone'}

        phi_elements = self.extract_phi_elements(filepath)
        found_types = {elem.element_type for elem in phi_elements}

        missing = expected_phi_elements - found_types
        return len(missing) == 0

    def check_phi_negative(self, filepath: str) -> bool:
        """
        Check if a PHI-negative document has NO patient identifiers

        Args:
            filepath: Path to document

        Returns:
            True if no PHI elements are found
        """
        phi_elements = self.extract_phi_elements(filepath)
        return len(phi_elements) == 0

    def check_file_integrity(self, filepath: str) -> bool:
        """
        Check if file can be opened and is not corrupted

        Args:
            filepath: Path to document

        Returns:
            True if file integrity is OK
        """
        filepath = Path(filepath)

        if not filepath.exists():
            return False

        if filepath.stat().st_size == 0:
            return False

        file_format = filepath.suffix.lower()

        try:
            if file_format == '.docx':
                return self._check_docx_integrity(filepath)
            elif file_format == '.pdf':
                return self._check_pdf_integrity(filepath)
            elif file_format == '.xlsx':
                return self._check_xlsx_integrity(filepath)
            elif file_format == '.eml':
                return self._check_eml_integrity(filepath)
            elif file_format == '.pptx':
                return self._check_pptx_integrity(filepath)
            else:
                return False
        except Exception:
            return False

    def extract_phi_elements(self, filepath: str) -> List[PHIElement]:
        """
        Extract all PHI elements from a document

        Args:
            filepath: Path to document

        Returns:
            List of PHIElement objects found in the document
        """
        filepath = Path(filepath)
        file_format = filepath.suffix.lower()

        if file_format == '.docx':
            return self._extract_phi_from_docx(filepath)
        elif file_format == '.pdf':
            return self._extract_phi_from_pdf(filepath)
        elif file_format == '.xlsx':
            return self._extract_phi_from_xlsx(filepath)
        elif file_format == '.eml':
            return self._extract_phi_from_eml(filepath)
        elif file_format == '.pptx':
            return self._extract_phi_from_pptx(filepath)
        else:
            return []

    def generate_validation_report(self, results: List[ValidationResult],
                                   output_path: str = None) -> Dict:
        """
        Generate a comprehensive validation report

        Args:
            results: List of ValidationResult objects
            output_path: Optional path to save JSON report

        Returns:
            Dictionary with validation statistics and details
        """
        total = len(results)
        passed = sum(1 for r in results if r.is_valid)
        failed = total - passed

        phi_positive_results = [r for r in results if r.expected_phi_type == 'positive']
        phi_negative_results = [r for r in results if r.expected_phi_type == 'negative']

        report = {
            'summary': {
                'total_documents': total,
                'passed': passed,
                'failed': failed,
                'pass_rate': round(passed / total * 100, 2) if total > 0 else 0,
                'phi_positive_count': len(phi_positive_results),
                'phi_negative_count': len(phi_negative_results),
            },
            'phi_positive_stats': self._generate_phi_type_stats(phi_positive_results),
            'phi_negative_stats': self._generate_phi_type_stats(phi_negative_results),
            'format_breakdown': self._generate_format_breakdown(results),
            'common_errors': self._analyze_common_errors(results),
            'phi_element_distribution': self._analyze_phi_distribution(results),
            'generated_at': datetime.now().isoformat(),
            'detailed_results': [r.to_dict() for r in results]
        }

        if output_path:
            output_path = Path(output_path)
            output_path.parent.mkdir(parents=True, exist_ok=True)
            with open(output_path, 'w') as f:
                json.dump(report, f, indent=2)

        return report

    # Private helper methods

    def _validate_phi_positive(self, phi_elements: List[PHIElement],
                               errors: List[str], warnings: List[str]) -> bool:
        """Validate PHI-positive document has required elements"""
        required = {'name', 'dob', 'mrn', 'ssn', 'address', 'phone'}
        found = {elem.element_type for elem in phi_elements}
        missing = required - found

        if missing:
            errors.append(f"Missing required PHI elements: {', '.join(missing)}")
            return False

        # Warn if some elements are found multiple times (might indicate duplicates)
        type_counts = {}
        for elem in phi_elements:
            type_counts[elem.element_type] = type_counts.get(elem.element_type, 0) + 1

        for elem_type, count in type_counts.items():
            if count > 3:
                warnings.append(f"Element type '{elem_type}' found {count} times - verify this is expected")

        return True

    def _validate_phi_negative(self, phi_elements: List[PHIElement],
                               errors: List[str], warnings: List[str]) -> bool:
        """Validate PHI-negative document has no patient identifiers"""
        if phi_elements:
            found_types = {elem.element_type for elem in phi_elements}
            errors.append(f"Found PHI elements in negative document: {', '.join(found_types)}")
            return False
        return True

    def _extract_metadata(self, filepath: Path) -> Dict:
        """Extract file metadata"""
        stat = filepath.stat()
        return {
            'filename': filepath.name,
            'file_size': stat.st_size,
            'created': datetime.fromtimestamp(stat.st_ctime).isoformat(),
            'modified': datetime.fromtimestamp(stat.st_mtime).isoformat(),
            'format': filepath.suffix.lower(),
        }

    # Format-specific integrity checks

    def _check_docx_integrity(self, filepath: Path) -> bool:
        """Check DOCX file integrity"""
        if not DocxDocument:
            return False
        try:
            # DOCX is a ZIP file
            if not zipfile.is_zipfile(filepath):
                return False
            doc = DocxDocument(str(filepath))
            # Try to access paragraphs to ensure document structure is valid
            _ = len(doc.paragraphs)
            return True
        except Exception:
            return False

    def _check_pdf_integrity(self, filepath: Path) -> bool:
        """Check PDF file integrity"""
        # Try with PyPDF2 first
        if PyPDF2:
            try:
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    _ = len(reader.pages)
                    return True
            except Exception:
                pass

        # Try with pdfplumber
        if pdfplumber:
            try:
                with pdfplumber.open(filepath) as pdf:
                    _ = len(pdf.pages)
                    return True
            except Exception:
                pass

        return False

    def _check_xlsx_integrity(self, filepath: Path) -> bool:
        """Check XLSX file integrity"""
        if not load_workbook:
            return False
        try:
            # XLSX is a ZIP file
            if not zipfile.is_zipfile(filepath):
                return False
            wb = load_workbook(filename=str(filepath), read_only=True)
            _ = wb.sheetnames
            wb.close()
            return True
        except Exception:
            return False

    def _check_eml_integrity(self, filepath: Path) -> bool:
        """Check EML file integrity"""
        if not extract_msg:
            # Fallback: try to read as text
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(100)
                    return len(content) > 0
            except Exception:
                return False

        try:
            msg = extract_msg.Message(str(filepath))
            _ = msg.subject
            return True
        except Exception:
            # Try as plain text email
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read(100)
                    return 'From:' in content or 'Subject:' in content
            except Exception:
                return False

    def _check_pptx_integrity(self, filepath: Path) -> bool:
        """Check PPTX file integrity"""
        if not Presentation:
            return False
        try:
            # PPTX is a ZIP file
            if not zipfile.is_zipfile(filepath):
                return False
            prs = Presentation(str(filepath))
            _ = len(prs.slides)
            return True
        except Exception:
            return False

    # Format-specific PHI extraction

    def _extract_phi_from_docx(self, filepath: Path) -> List[PHIElement]:
        """Extract PHI from DOCX files"""
        if not DocxDocument:
            return []

        phi_elements = []
        try:
            doc = DocxDocument(str(filepath))

            # Extract from paragraphs
            for i, para in enumerate(doc.paragraphs):
                text = para.text
                location = f"paragraph_{i+1}"
                phi_elements.extend(self._find_phi_in_text(text, location))

            # Extract from tables
            for table_idx, table in enumerate(doc.tables):
                for row_idx, row in enumerate(table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        text = cell.text
                        location = f"table_{table_idx+1}_row_{row_idx+1}_cell_{cell_idx+1}"
                        phi_elements.extend(self._find_phi_in_text(text, location))

        except Exception:
            pass

        return self._deduplicate_phi_elements(phi_elements)

    def _extract_phi_from_pdf(self, filepath: Path) -> List[PHIElement]:
        """Extract PHI from PDF files"""
        phi_elements = []

        # Try pdfplumber first (better text extraction)
        if pdfplumber:
            try:
                with pdfplumber.open(filepath) as pdf:
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            location = f"page_{page_num}"
                            phi_elements.extend(self._find_phi_in_text(text, location))
                return self._deduplicate_phi_elements(phi_elements)
            except Exception:
                pass

        # Fallback to PyPDF2
        if PyPDF2:
            try:
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page_num, page in enumerate(reader.pages, 1):
                        text = page.extract_text()
                        if text:
                            location = f"page_{page_num}"
                            phi_elements.extend(self._find_phi_in_text(text, location))
            except Exception:
                pass

        return self._deduplicate_phi_elements(phi_elements)

    def _extract_phi_from_xlsx(self, filepath: Path) -> List[PHIElement]:
        """Extract PHI from XLSX files"""
        if not load_workbook:
            return []

        phi_elements = []
        try:
            wb = load_workbook(filename=str(filepath), read_only=True, data_only=True)

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), 1):
                    for col_idx, cell_value in enumerate(row, 1):
                        if cell_value:
                            text = str(cell_value)
                            location = f"sheet_{sheet_name}_row_{row_idx}_col_{col_idx}"
                            phi_elements.extend(self._find_phi_in_text(text, location))

            wb.close()
        except Exception:
            pass

        return self._deduplicate_phi_elements(phi_elements)

    def _extract_phi_from_eml(self, filepath: Path) -> List[PHIElement]:
        """Extract PHI from EML files"""
        phi_elements = []

        if extract_msg:
            try:
                msg = extract_msg.Message(str(filepath))

                # Check subject
                if msg.subject:
                    phi_elements.extend(self._find_phi_in_text(msg.subject, "subject"))

                # Check body
                if msg.body:
                    phi_elements.extend(self._find_phi_in_text(msg.body, "body"))

                return self._deduplicate_phi_elements(phi_elements)
            except Exception:
                pass

        # Fallback: parse as plain text
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
                phi_elements.extend(self._find_phi_in_text(content, "email_content"))
        except Exception:
            pass

        return self._deduplicate_phi_elements(phi_elements)

    def _extract_phi_from_pptx(self, filepath: Path) -> List[PHIElement]:
        """Extract PHI from PPTX files"""
        if not Presentation:
            return []

        phi_elements = []
        try:
            prs = Presentation(str(filepath))

            for slide_idx, slide in enumerate(prs.slides, 1):
                # Extract from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        location = f"slide_{slide_idx}"
                        phi_elements.extend(self._find_phi_in_text(shape.text, location))

                    # Extract from tables in slides
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row_idx, row in enumerate(table.rows):
                            for cell_idx, cell in enumerate(row.cells):
                                if cell.text:
                                    location = f"slide_{slide_idx}_table_row_{row_idx+1}_cell_{cell_idx+1}"
                                    phi_elements.extend(self._find_phi_in_text(cell.text, location))

        except Exception:
            pass

        return self._deduplicate_phi_elements(phi_elements)

    def _find_phi_in_text(self, text: str, location: str) -> List[PHIElement]:
        """Find all PHI elements in a text string"""
        phi_elements = []

        # SSN
        for match in self.patterns.SSN.finditer(text):
            phi_elements.append(PHIElement('ssn', match.group(), location))

        # MRN
        for match in self.patterns.MRN.finditer(text):
            phi_elements.append(PHIElement('mrn', match.group(), location))

        # Phone
        for match in self.patterns.PHONE.finditer(text):
            phi_elements.append(PHIElement('phone', match.group(), location))
        for match in self.patterns.PHONE_ALT.finditer(text):
            # Avoid duplicates if already found with primary pattern
            phone = match.group()
            if not any(elem.value == phone for elem in phi_elements if elem.element_type == 'phone'):
                phi_elements.append(PHIElement('phone', phone, location))

        # DOB
        for match in self.patterns.DOB.finditer(text):
            phi_elements.append(PHIElement('dob', match.group(), location))

        # Email
        for match in self.patterns.EMAIL.finditer(text):
            phi_elements.append(PHIElement('email', match.group(), location))

        # Address
        for match in self.patterns.ADDRESS.finditer(text):
            phi_elements.append(PHIElement('address', match.group(), location))

        # ZIP code (part of address)
        for match in self.patterns.ZIP_CODE.finditer(text):
            # Only add if we haven't found an address yet in this location
            if not any(elem.element_type == 'address' and elem.location == location for elem in phi_elements):
                phi_elements.append(PHIElement('address', match.group(), location))

        # Name
        for match in self.patterns.NAME.finditer(text):
            name = match.group(1)
            phi_elements.append(PHIElement('name', name, location))

        # Also look for names in common patterns
        # "Dr. FirstName LastName" or just "FirstName LastName" at start of line
        name_pattern = re.compile(r'\b(?:Dr\.\s+)?([A-Z][a-z]+\s+[A-Z][a-z]+)\b')
        for match in name_pattern.finditer(text):
            name = match.group(1)
            # Only add if name is not already found
            if not any(elem.value == name for elem in phi_elements if elem.element_type == 'name'):
                phi_elements.append(PHIElement('name', name, location, confidence=0.7))

        # Insurance ID
        for match in self.patterns.INSURANCE_ID.finditer(text):
            phi_elements.append(PHIElement('insurance_id', match.group(), location))

        return phi_elements

    def _deduplicate_phi_elements(self, phi_elements: List[PHIElement]) -> List[PHIElement]:
        """Remove duplicate PHI elements"""
        seen = set()
        unique = []

        for elem in phi_elements:
            key = (elem.element_type, elem.value)
            if key not in seen:
                seen.add(key)
                unique.append(elem)

        return unique

    def _generate_phi_type_stats(self, results: List[ValidationResult]) -> Dict:
        """Generate statistics for PHI-positive or PHI-negative results"""
        if not results:
            return {
                'count': 0,
                'passed': 0,
                'failed': 0,
                'pass_rate': 0
            }

        passed = sum(1 for r in results if r.is_valid)
        return {
            'count': len(results),
            'passed': passed,
            'failed': len(results) - passed,
            'pass_rate': round(passed / len(results) * 100, 2)
        }

    def _generate_format_breakdown(self, results: List[ValidationResult]) -> Dict:
        """Generate breakdown by file format"""
        breakdown = {}

        for result in results:
            fmt = result.file_format
            if fmt not in breakdown:
                breakdown[fmt] = {'total': 0, 'passed': 0, 'failed': 0}

            breakdown[fmt]['total'] += 1
            if result.is_valid:
                breakdown[fmt]['passed'] += 1
            else:
                breakdown[fmt]['failed'] += 1

        # Add pass rates
        for fmt in breakdown:
            total = breakdown[fmt]['total']
            if total > 0:
                breakdown[fmt]['pass_rate'] = round(breakdown[fmt]['passed'] / total * 100, 2)

        return breakdown

    def _analyze_common_errors(self, results: List[ValidationResult]) -> List[Dict]:
        """Analyze and return most common errors"""
        error_counts = {}

        for result in results:
            for error in result.validation_errors:
                error_counts[error] = error_counts.get(error, 0) + 1

        # Sort by frequency
        sorted_errors = sorted(error_counts.items(), key=lambda x: x[1], reverse=True)

        return [
            {'error': error, 'count': count}
            for error, count in sorted_errors[:10]  # Top 10 errors
        ]

    def _analyze_phi_distribution(self, results: List[ValidationResult]) -> Dict:
        """Analyze distribution of PHI element types found"""
        element_counts = {}

        for result in results:
            for elem in result.phi_elements_found:
                elem_type = elem.element_type
                element_counts[elem_type] = element_counts.get(elem_type, 0) + 1

        return element_counts


def validate_batch(file_paths: List[str], expected_phi_type: str = 'unknown',
                   output_report: str = None) -> Dict:
    """
    Convenience function to validate multiple files and generate a report

    Args:
        file_paths: List of file paths to validate
        expected_phi_type: 'positive', 'negative', or 'unknown'
        output_report: Optional path to save JSON report

    Returns:
        Validation report dictionary
    """
    validator = PHIValidator()
    results = []

    for filepath in file_paths:
        result = validator.validate_document(filepath, expected_phi_type)
        results.append(result)

    report = validator.generate_validation_report(results, output_report)
    return report


# Example usage
if __name__ == '__main__':
    # Example validation
    validator = PHIValidator()

    # Validate a single document
    result = validator.validate_document(
        '/path/to/document.docx',
        expected_phi_type='positive'
    )

    print(f"Valid: {result.is_valid}")
    print(f"PHI Elements Found: {len(result.phi_elements_found)}")
    for elem in result.phi_elements_found:
        print(f"  - {elem.element_type}: {elem.value} (at {elem.location})")

    # Batch validation
    files = [
        '/path/to/doc1.docx',
        '/path/to/doc2.pdf',
        '/path/to/doc3.xlsx',
    ]

    report = validate_batch(
        files,
        expected_phi_type='positive',
        output_report='/path/to/validation_report.json'
    )

    print("\nValidation Report:")
    print(f"Total: {report['summary']['total_documents']}")
    print(f"Passed: {report['summary']['passed']}")
    print(f"Failed: {report['summary']['failed']}")
    print(f"Pass Rate: {report['summary']['pass_rate']}%")
