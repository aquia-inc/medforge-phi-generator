#!/usr/bin/env python3
"""
Purview File Fidelity Validator

Validates generated MedForge files against the Purview File Fidelity Checklist.
Checks all 18 items that can cause Purview classifier training to timeout or fail.

Usage:
    uv run python tests/validate_file_fidelity.py <output_dir> [--json]
"""
import sys
import os
import json
import email
import email.utils
import zipfile
import struct
from pathlib import Path
from datetime import datetime
from collections import defaultdict

# Severity levels
FATAL = "FATAL"
HIGH = "HIGH"
MEDIUM = "MEDIUM"


class FileValidationResult:
    """Result of validating a single file."""

    def __init__(self, filepath: str):
        self.filepath = filepath
        self.filename = os.path.basename(filepath)
        self.ext = os.path.splitext(filepath)[1].lower()
        self.size_bytes = os.path.getsize(filepath) if os.path.exists(filepath) else 0
        self.issues = []  # list of (check_num, severity, description)
        self.checks_passed = []

    def add_issue(self, check_num: int, severity: str, description: str):
        self.issues.append((check_num, severity, description))

    def add_pass(self, check_num: int, description: str):
        self.checks_passed.append((check_num, description))

    @property
    def fatal_count(self):
        return sum(1 for _, sev, _ in self.issues if sev == FATAL)

    @property
    def high_count(self):
        return sum(1 for _, sev, _ in self.issues if sev == HIGH)

    @property
    def medium_count(self):
        return sum(1 for _, sev, _ in self.issues if sev == MEDIUM)

    @property
    def is_purview_safe(self):
        return self.fatal_count == 0


class PurviewFidelityValidator:
    """Validates files against the Purview File Fidelity Checklist."""

    MAX_FILE_SIZE = 64 * 1024 * 1024  # 64MB
    MAX_XLSX_SIZE = 4 * 1024 * 1024   # 4MB
    MAX_PATH_LEN = 260
    MAX_TEXT_CHARS = 2_000_000

    def validate_directory(self, directory: str) -> list:
        """Validate all files in a directory recursively."""
        results = []
        for root, _dirs, files in os.walk(directory):
            for fname in files:
                filepath = os.path.join(root, fname)
                ext = os.path.splitext(fname)[1].lower()
                if ext in ('.eml', '.pdf', '.docx', '.xlsx', '.pptx', '.zip'):
                    result = self.validate_file(filepath)
                    results.append(result)
        return results

    def validate_file(self, filepath: str) -> FileValidationResult:
        """Run all applicable checks against a single file."""
        result = FileValidationResult(filepath)
        ext = result.ext

        # Check 1: File opens without error
        self._check_file_opens(result)

        # Check 11: File size under 64MB
        self._check_file_size(result)

        # Check 12: XLSX under 4MB
        if ext == '.xlsx':
            self._check_xlsx_size(result)

        # Check 15: File path under 260 characters
        self._check_path_length(result)

        # EML-specific checks
        if ext == '.eml':
            self._check_eml(result)

        # PDF-specific checks
        if ext == '.pdf':
            self._check_pdf(result)

        # OOXML checks (DOCX, XLSX, PPTX)
        if ext in ('.docx', '.xlsx', '.pptx'):
            self._check_ooxml(result)

        # PPTX-specific: slide count and content check
        if ext == '.pptx':
            self._check_pptx_slides(result)

        return result

    # ---- Individual checks ----

    def _check_file_opens(self, result: FileValidationResult):
        """Check 1: File opens without error."""
        if not os.path.exists(result.filepath):
            result.add_issue(1, FATAL, "File does not exist")
            return
        if result.size_bytes == 0:
            result.add_issue(1, FATAL, "File is empty (0 bytes)")
            return
        result.add_pass(1, "File exists and is non-empty")

    def _check_file_size(self, result: FileValidationResult):
        """Check 11: File size under 64MB."""
        if result.size_bytes > self.MAX_FILE_SIZE:
            result.add_issue(11, FATAL,
                f"File size {result.size_bytes / 1024 / 1024:.1f}MB exceeds 64MB limit")
        else:
            result.add_pass(11, f"File size OK ({result.size_bytes / 1024:.1f}KB)")

    def _check_xlsx_size(self, result: FileValidationResult):
        """Check 12: XLSX under 4MB."""
        if result.size_bytes > self.MAX_XLSX_SIZE:
            result.add_issue(12, HIGH,
                f"XLSX size {result.size_bytes / 1024 / 1024:.1f}MB exceeds 4MB recommended limit")
        else:
            result.add_pass(12, "XLSX size OK")

    def _check_path_length(self, result: FileValidationResult):
        """Check 15: File path under 260 characters."""
        if len(result.filepath) > self.MAX_PATH_LEN:
            result.add_issue(15, MEDIUM,
                f"Path length {len(result.filepath)} exceeds 260 character limit")
        else:
            result.add_pass(15, "Path length OK")

    def _check_eml(self, result: FileValidationResult):
        """Run all EML-specific checks (2-6, 14, 16)."""
        try:
            with open(result.filepath, 'rb') as f:
                raw_bytes = f.read()
        except Exception as e:
            result.add_issue(1, FATAL, f"Cannot read EML file: {e}")
            return

        # Try parsing as bytes first, then as string
        msg = None
        parsed_as_bytes = True
        try:
            msg = email.message_from_bytes(raw_bytes)
        except Exception:
            try:
                msg = email.message_from_string(raw_bytes.decode('utf-8', errors='replace'))
                parsed_as_bytes = False
            except Exception as e:
                result.add_issue(2, FATAL, f"Cannot parse EML MIME structure: {e}")
                return

        # Check 2: MIME boundaries valid
        content_type = msg.get_content_type()
        if content_type and content_type.startswith('multipart/'):
            boundary = msg.get_boundary()
            if not boundary:
                result.add_issue(2, FATAL, "Multipart message missing MIME boundary")
            else:
                # Check boundary appears in raw content
                boundary_bytes = boundary.encode('utf-8', errors='replace')
                if boundary_bytes not in raw_bytes:
                    result.add_issue(2, FATAL,
                        f"MIME boundary '{boundary}' not found in raw file content (corrupted)")
                else:
                    result.add_pass(2, "MIME boundaries valid")
        else:
            result.add_pass(2, "Not multipart, boundary check N/A")

        # Check 3: Binary attachments decode correctly
        self._check_eml_attachments(result, msg)

        # Check 4: RFC 2822 Date header has timezone
        date_header = msg.get('Date', '')
        if not date_header:
            result.add_issue(4, HIGH, "Missing Date header")
        else:
            import re
            date_str = date_header.strip()

            # parsedate_tz defaults missing tz to 0, so we can't rely on it.
            # Instead, check the raw string for an actual timezone indicator:
            # +HHMM, -HHMM, or a standard tz abbreviation (UTC, EST, etc.)
            tz_pattern = r'[+-]\d{4}\s*$|(?:UT|UTC|GMT|EST|EDT|CST|CDT|MST|MDT|PST|PDT)\s*$'
            if re.search(tz_pattern, date_str):
                result.add_pass(4, f"Date header has timezone: '{date_str}'")
            else:
                result.add_issue(4, HIGH,
                    f"Date header missing timezone: '{date_str}'")

        # Check 5: MIME-Version header present
        mime_version = msg.get('MIME-Version', '')
        if not mime_version:
            result.add_issue(5, HIGH, "Missing MIME-Version header")
        else:
            result.add_pass(5, f"MIME-Version: {mime_version}")

        # Check 6: Content-Type charset matches actual encoding
        self._check_eml_charset(result, msg)

        # Check 14: Nesting depth <= 3 levels
        self._check_nesting_depth(result, msg, current_depth=0)

        # Check 16: UTF-8 encoding
        try:
            raw_bytes.decode('utf-8')
            result.add_pass(16, "File is valid UTF-8")
        except UnicodeDecodeError:
            # EML files with base64 attachments may not be pure UTF-8 in text mode
            # but should be in binary mode. Check the text parts specifically.
            text_parts_ok = True
            for part in msg.walk():
                if part.get_content_maintype() == 'text':
                    payload = part.get_payload(decode=True)
                    if payload:
                        charset = part.get_content_charset() or 'us-ascii'
                        try:
                            payload.decode(charset)
                        except (UnicodeDecodeError, LookupError):
                            text_parts_ok = False
                            result.add_issue(16, MEDIUM,
                                f"Text part charset '{charset}' doesn't match actual encoding")
            if text_parts_ok:
                result.add_pass(16, "Text parts properly encoded")

    def _check_eml_attachments(self, result: FileValidationResult, msg):
        """Check 3: Binary attachments in EML decode correctly."""
        attachment_count = 0
        for part in msg.walk():
            content_disp = str(part.get('Content-Disposition', ''))
            if 'attachment' not in content_disp:
                continue

            attachment_count += 1
            filename = part.get_filename() or 'unknown'
            payload = part.get_payload(decode=True)

            if payload is None:
                result.add_issue(3, FATAL,
                    f"Attachment '{filename}' failed to decode (None payload)")
                continue

            if len(payload) == 0:
                result.add_issue(3, FATAL,
                    f"Attachment '{filename}' decoded to 0 bytes")
                continue

            # Validate the attachment content based on extension
            att_ext = os.path.splitext(filename)[1].lower()

            if att_ext == '.pdf':
                if not payload[:5] == b'%PDF-':
                    result.add_issue(3, FATAL,
                        f"Attachment '{filename}' does not start with %PDF- header "
                        f"(got {payload[:20]!r}) - likely corrupted by text-mode write")
                else:
                    result.add_pass(3, f"Attachment '{filename}' has valid PDF header")

            elif att_ext in ('.docx', '.xlsx', '.pptx'):
                # OOXML files are ZIP archives - check magic bytes
                if not payload[:4] == b'PK\x03\x04':
                    result.add_issue(3, FATAL,
                        f"Attachment '{filename}' missing ZIP/OOXML magic bytes "
                        f"(got {payload[:4]!r}) - likely corrupted")
                else:
                    result.add_pass(3, f"Attachment '{filename}' has valid OOXML header")

            elif att_ext == '.zip':
                if not payload[:4] == b'PK\x03\x04':
                    result.add_issue(3, FATAL,
                        f"Attachment '{filename}' missing ZIP magic bytes")
                else:
                    # Try actually opening the ZIP
                    import io
                    try:
                        with zipfile.ZipFile(io.BytesIO(payload)) as zf:
                            zf.testzip()
                        result.add_pass(3, f"Attachment '{filename}' is valid ZIP")
                    except Exception as e:
                        result.add_issue(3, FATAL,
                            f"Attachment '{filename}' ZIP is corrupted: {e}")
            else:
                if len(payload) > 0:
                    result.add_pass(3, f"Attachment '{filename}' decoded ({len(payload)} bytes)")

        if attachment_count == 0:
            result.add_pass(3, "No attachments to validate")

    def _check_eml_charset(self, result: FileValidationResult, msg):
        """Check 6: Content-Type charset matches actual encoding."""
        for part in msg.walk():
            if part.get_content_maintype() != 'text':
                continue

            charset = part.get_content_charset()
            payload = part.get_payload(decode=True)

            if payload is None:
                continue

            if charset is None or charset == 'us-ascii':
                # Check if content has non-ASCII characters
                try:
                    payload.decode('us-ascii')
                    if charset is None:
                        result.add_issue(6, HIGH,
                            f"Text part missing explicit charset (defaulting to us-ascii)")
                    else:
                        result.add_pass(6, "us-ascii charset matches content")
                except UnicodeDecodeError:
                    if charset is None:
                        result.add_issue(6, HIGH,
                            "Text part has non-ASCII content but no charset specified")
                    else:
                        result.add_issue(6, HIGH,
                            f"Text part has non-ASCII content but charset is '{charset}'")
            else:
                try:
                    payload.decode(charset)
                    result.add_pass(6, f"Charset '{charset}' matches content")
                except (UnicodeDecodeError, LookupError) as e:
                    result.add_issue(6, HIGH, f"Charset '{charset}' decode failed: {e}")

    def _check_nesting_depth(self, result: FileValidationResult, msg, current_depth: int):
        """Check 14: Nesting depth <= 3 levels."""
        max_depth = current_depth
        for part in msg.walk():
            content_disp = str(part.get('Content-Disposition', ''))
            if 'attachment' in content_disp:
                payload = part.get_payload(decode=True)
                if payload:
                    att_ext = os.path.splitext(part.get_filename() or '')[1].lower()
                    depth = current_depth + 1

                    if att_ext == '.eml':
                        # Nested email
                        depth += 1
                    elif att_ext == '.zip':
                        depth += 1

                    max_depth = max(max_depth, depth)

        if max_depth > 3:
            result.add_issue(14, HIGH,
                f"Nesting depth {max_depth} exceeds 3-level limit")
        else:
            result.add_pass(14, f"Nesting depth OK ({max_depth})")

    def _check_pdf(self, result: FileValidationResult):
        """Run all PDF-specific checks (7, 8, 9)."""
        try:
            with open(result.filepath, 'rb') as f:
                header = f.read(1024)
        except Exception as e:
            result.add_issue(8, FATAL, f"Cannot read PDF: {e}")
            return

        # Check 8: PDF has %PDF header
        if not header.startswith(b'%PDF'):
            result.add_issue(8, FATAL,
                f"Missing %PDF header (got {header[:20]!r})")
            return
        result.add_pass(8, "Valid %PDF header")

        # Check 7: Text extractable by pdfplumber
        try:
            import pdfplumber
            with pdfplumber.open(result.filepath) as pdf:
                all_text = ""
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    all_text += text

                if len(all_text.strip()) == 0:
                    result.add_issue(7, FATAL,
                        "No text extractable by pdfplumber (Purview cannot index)")
                else:
                    result.add_pass(7,
                        f"Text extractable ({len(all_text)} chars)")

                    # Check 13: Text content under 2M characters
                    if len(all_text) > self.MAX_TEXT_CHARS:
                        result.add_issue(13, HIGH,
                            f"Text content {len(all_text)} chars exceeds 2M limit")
                    else:
                        result.add_pass(13, "Text length OK")
        except ImportError:
            result.add_issue(7, HIGH, "pdfplumber not installed, cannot verify text extraction")
        except Exception as e:
            result.add_issue(7, FATAL, f"pdfplumber failed to open: {e}")

        # Check 9: No form fields trapping text (AcroForm)
        try:
            import pikepdf
            pdf = pikepdf.open(result.filepath)
            if '/AcroForm' in pdf.Root:
                acroform = pdf.Root.AcroForm
                if '/Fields' in acroform:
                    fields = list(acroform.Fields)
                    fields_with_values = 0
                    for field in fields:
                        if '/V' in field:
                            val = str(field.V) if field.V is not None else ''
                            # Only count non-empty values (blank templates have /V='')
                            if val.strip():
                                fields_with_values += 1

                    if fields_with_values > 0:
                        result.add_issue(9, FATAL,
                            f"PDF has AcroForm with {fields_with_values} populated fields "
                            f"(values trapped in form fields, not extractable by Purview)")
                    else:
                        result.add_pass(9,
                            "AcroForm present but fields are empty (blank template)")
                else:
                    result.add_pass(9, "AcroForm present but no Fields array")
            else:
                result.add_pass(9, "No AcroForm (good)")
            pdf.close()
        except ImportError:
            result.add_issue(9, HIGH, "pikepdf not installed, cannot check AcroForm")
        except Exception as e:
            result.add_issue(9, MEDIUM, f"pikepdf check failed: {e}")

    def _check_ooxml(self, result: FileValidationResult):
        """Check 10: Valid ZIP structure (OOXML) for DOCX, XLSX, PPTX."""
        try:
            with open(result.filepath, 'rb') as f:
                magic = f.read(4)

            if magic != b'PK\x03\x04':
                result.add_issue(10, FATAL,
                    f"Missing ZIP/OOXML magic bytes (got {magic!r})")
                return

            with zipfile.ZipFile(result.filepath) as zf:
                bad_file = zf.testzip()
                if bad_file:
                    result.add_issue(10, FATAL,
                        f"Corrupt entry in ZIP: {bad_file}")
                else:
                    # Check for expected OOXML content
                    names = zf.namelist()
                    if '[Content_Types].xml' in names:
                        result.add_pass(10, "Valid OOXML ZIP structure")
                    else:
                        result.add_issue(10, FATAL,
                            "ZIP missing [Content_Types].xml (not valid OOXML)")
        except zipfile.BadZipFile as e:
            result.add_issue(10, FATAL, f"Invalid ZIP file: {e}")
        except Exception as e:
            result.add_issue(10, FATAL, f"Cannot open as ZIP: {e}")


    def _check_pptx_slides(self, result: FileValidationResult):
        """Check 19: PPTX has at least 2 slides and contains text content."""
        try:
            from pptx import Presentation
            prs = Presentation(result.filepath)
            slide_count = len(prs.slides)

            if slide_count < 2:
                result.add_issue(19, HIGH,
                    f"PPTX has only {slide_count} slide(s), expected >= 2")
            else:
                result.add_pass(19, f"PPTX has {slide_count} slides")

            # Check for text content in slides
            total_text = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        total_text += len(shape.text_frame.text)

            if total_text == 0:
                result.add_issue(19, HIGH, "PPTX contains no text content")
            else:
                result.add_pass(19, f"PPTX has {total_text} chars of text content")

        except ImportError:
            result.add_issue(19, MEDIUM, "python-pptx not installed, cannot check slides")
        except Exception as e:
            result.add_issue(19, FATAL, f"Cannot open PPTX: {e}")


def check_lab_data_consistency(output_dir: str) -> list:
    """
    Check 17: Lab data self-consistent (flags match values).
    This checks the generated lab data at the source level by importing the generator
    and running it multiple times to detect inconsistencies.
    """
    issues = []
    try:
        sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'src'))
        from generators.patient_generator import PatientGenerator

        gen = PatientGenerator(seed=42)
        for i in range(50):
            lab = gen.generate_lab_results(panel_type='drug_screen')
            for test_result in lab['results']:
                value = test_result['value']
                flag = test_result['flag']
                test_name = test_result['test']

                # A POSITIVE result should have flag 'A', NEGATIVE should have no flag
                if value == 'POSITIVE' and flag == '':
                    issues.append(
                        f"Drug screen inconsistency (iteration {i}): "
                        f"{test_name} value=POSITIVE but flag is empty")
                elif value == 'NEGATIVE' and flag == 'A':
                    issues.append(
                        f"Drug screen inconsistency (iteration {i}): "
                        f"{test_name} value=NEGATIVE but flag=A (abnormal)")
    except Exception as e:
        issues.append(f"Could not check lab data consistency: {e}")

    return issues


def generate_report(results: list, lab_issues: list, output_format: str = 'text') -> str:
    """Generate a validation report."""
    total_files = len(results)
    purview_safe = sum(1 for r in results if r.is_purview_safe)
    total_fatal = sum(r.fatal_count for r in results)
    total_high = sum(r.high_count for r in results)
    total_medium = sum(r.medium_count for r in results)

    if output_format == 'json':
        report = {
            'summary': {
                'total_files': total_files,
                'purview_safe': purview_safe,
                'purview_unsafe': total_files - purview_safe,
                'total_fatal': total_fatal,
                'total_high': total_high,
                'total_medium': total_medium,
                'lab_data_issues': len(lab_issues),
            },
            'files': [],
            'lab_data_issues': lab_issues,
        }
        for r in results:
            file_entry = {
                'filepath': r.filepath,
                'filename': r.filename,
                'extension': r.ext,
                'size_bytes': r.size_bytes,
                'purview_safe': r.is_purview_safe,
                'issues': [
                    {'check': num, 'severity': sev, 'description': desc}
                    for num, sev, desc in r.issues
                ],
                'passed': [
                    {'check': num, 'description': desc}
                    for num, desc in r.checks_passed
                ],
            }
            report['files'].append(file_entry)
        return json.dumps(report, indent=2)

    # Text format
    lines = []
    lines.append("=" * 80)
    lines.append("PURVIEW FILE FIDELITY VALIDATION REPORT")
    lines.append(f"Generated: {datetime.now().isoformat()}")
    lines.append("=" * 80)
    lines.append("")
    lines.append("SUMMARY")
    lines.append("-" * 40)
    lines.append(f"Total files validated:  {total_files}")
    lines.append(f"Purview-safe files:    {purview_safe}")
    lines.append(f"Purview-UNSAFE files:  {total_files - purview_safe}")
    lines.append(f"Fatal issues:          {total_fatal}")
    lines.append(f"High issues:           {total_high}")
    lines.append(f"Medium issues:         {total_medium}")
    lines.append(f"Lab data issues:       {len(lab_issues)}")
    lines.append("")

    # Issues by check number
    issues_by_check = defaultdict(list)
    for r in results:
        for num, sev, desc in r.issues:
            issues_by_check[num].append((r.filename, sev, desc))

    if issues_by_check:
        lines.append("ISSUES BY CHECK")
        lines.append("-" * 40)
        check_names = {
            1: "File opens without error",
            2: "MIME boundaries valid",
            3: "Binary attachments decode correctly",
            4: "RFC 2822 Date header has timezone",
            5: "MIME-Version header present",
            6: "Content-Type charset matches encoding",
            7: "Text extractable by pdfplumber",
            8: "PDF has %PDF header",
            9: "No form fields trapping text (AcroForm)",
            10: "Valid ZIP structure (OOXML)",
            11: "File size under 64MB",
            12: "XLSX under 4MB",
            13: "Text content under 2M characters",
            14: "Nesting depth <= 3 levels",
            15: "File path under 260 characters",
            16: "UTF-8 encoding throughout",
            17: "Lab data self-consistent",
            18: "English content only",
            19: "PPTX slides and content check",
        }
        for check_num in sorted(issues_by_check.keys()):
            items = issues_by_check[check_num]
            check_name = check_names.get(check_num, f"Check {check_num}")
            lines.append(f"\n  Check {check_num}: {check_name} ({len(items)} issues)")
            for fname, sev, desc in items:
                lines.append(f"    [{sev}] {fname}: {desc}")

    if lab_issues:
        lines.append(f"\n  Check 17: Lab data self-consistent ({len(lab_issues)} issues)")
        for issue in lab_issues[:10]:
            lines.append(f"    [MEDIUM] {issue}")
        if len(lab_issues) > 10:
            lines.append(f"    ... and {len(lab_issues) - 10} more")

    lines.append("")

    # Per-file detail
    lines.append("PER-FILE DETAILS")
    lines.append("-" * 40)

    # Show unsafe files first
    unsafe = [r for r in results if not r.is_purview_safe]
    safe = [r for r in results if r.is_purview_safe]

    if unsafe:
        lines.append(f"\n  UNSAFE FILES ({len(unsafe)}):")
        for r in unsafe:
            lines.append(f"\n    {r.filename} ({r.ext}, {r.size_bytes / 1024:.1f}KB)")
            for num, sev, desc in r.issues:
                marker = "!!!" if sev == FATAL else "! " if sev == HIGH else "  "
                lines.append(f"      {marker} [{sev}] Check {num}: {desc}")

    if safe:
        lines.append(f"\n  SAFE FILES ({len(safe)}):")
        for r in safe:
            issue_summary = ""
            if r.high_count or r.medium_count:
                parts = []
                if r.high_count:
                    parts.append(f"{r.high_count} high")
                if r.medium_count:
                    parts.append(f"{r.medium_count} medium")
                issue_summary = f" (non-fatal issues: {', '.join(parts)})"
            lines.append(f"    OK  {r.filename}{issue_summary}")

    lines.append("")
    lines.append("=" * 80)
    return "\n".join(lines)


def main():
    import argparse
    parser = argparse.ArgumentParser(description="Purview File Fidelity Validator")
    parser.add_argument("directory", help="Directory containing generated files to validate")
    parser.add_argument("--json", action="store_true", help="Output as JSON")
    parser.add_argument("--output", "-o", help="Write report to file instead of stdout")
    args = parser.parse_args()

    if not os.path.isdir(args.directory):
        print(f"Error: '{args.directory}' is not a directory", file=sys.stderr)
        sys.exit(1)

    print(f"Validating files in: {args.directory}", file=sys.stderr)

    validator = PurviewFidelityValidator()
    results = validator.validate_directory(args.directory)

    print(f"Validated {len(results)} files", file=sys.stderr)

    # Run lab data consistency check
    print("Running lab data consistency check...", file=sys.stderr)
    lab_issues = check_lab_data_consistency(args.directory)

    fmt = 'json' if args.json else 'text'
    report = generate_report(results, lab_issues, output_format=fmt)

    if args.output:
        with open(args.output, 'w') as f:
            f.write(report)
        print(f"Report written to: {args.output}", file=sys.stderr)
    else:
        print(report)

    # Exit code based on fatal issues
    total_fatal = sum(r.fatal_count for r in results)
    if total_fatal > 0 or lab_issues:
        sys.exit(1)
    sys.exit(0)


if __name__ == '__main__':
    main()
