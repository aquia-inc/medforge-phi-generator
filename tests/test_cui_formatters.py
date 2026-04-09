"""
Unit tests for CUI document formatters.

Tests that each CUI formatter produces valid, openable output files.
Distinct from test_cui_generators.py which tests the data generation layer.
"""
import email
import os
import random

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


class TestCUIDocxFormatter:
    """Tests for CUI DOCX document generation."""

    def test_creates_valid_docx(self, tmp_output_dir, sample_cui_data):
        """create_cui_document produces an openable DOCX file."""
        from formatters.cui_formatter import CUIDocxFormatter

        fmt = CUIDocxFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['legal']
        filepath = fmt.create_cui_document(doc_data, "test_legal.docx")

        assert os.path.exists(filepath)
        assert filepath.endswith('.docx')
        doc = Document(filepath)
        assert len(doc.paragraphs) > 0

    def test_with_component_config(self, tmp_output_dir, sample_cui_data):
        """create_cui_document applies component configuration styling."""
        from formatters.cui_formatter import CUIDocxFormatter
        from templates.components import ComponentMixer

        mixer = ComponentMixer(seed=42)
        config = mixer.get_random_configuration()

        fmt = CUIDocxFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['financial']
        filepath = fmt.create_cui_document(doc_data, "test_styled.docx", component_config=config)

        assert os.path.exists(filepath)
        doc = Document(filepath)
        assert len(doc.paragraphs) > 0

    def test_contains_category_content(self, tmp_output_dir, sample_cui_data):
        """Generated DOCX for a category contains relevant text."""
        from formatters.cui_formatter import CUIDocxFormatter

        fmt = CUIDocxFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['tax']
        filepath = fmt.create_cui_document(doc_data, "test_tax.docx")

        doc = Document(filepath)
        full_text = '\n'.join(p.text for p in doc.paragraphs)
        # Tax documents should contain at least the title or some tax-related content
        assert len(full_text) > 50


class TestCUIPdfFormatter:
    """Tests for CUI PDF document generation."""

    def test_creates_valid_pdf(self, tmp_output_dir, sample_cui_data):
        """create_cui_pdf produces a valid PDF file."""
        from formatters.cui_formatter import CUIPdfFormatter

        fmt = CUIPdfFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['procurement']
        filepath = fmt.create_cui_pdf(doc_data, "test_procurement.pdf")

        assert os.path.exists(filepath)
        assert filepath.endswith('.pdf')
        # Verify it starts with PDF magic bytes
        with open(filepath, 'rb') as f:
            header = f.read(5)
        assert header == b'%PDF-'

    def test_with_component_config(self, tmp_output_dir, sample_cui_data):
        """create_cui_pdf applies component configuration (font mapping)."""
        from formatters.cui_formatter import CUIPdfFormatter
        from templates.components import ComponentMixer

        mixer = ComponentMixer(seed=42)
        config = mixer.get_random_configuration()

        fmt = CUIPdfFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['legal']
        filepath = fmt.create_cui_pdf(doc_data, "test_styled.pdf", component_config=config)

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            header = f.read(5)
        assert header == b'%PDF-'


class TestCUIXlsxFormatter:
    """Tests for CUI XLSX spreadsheet generation."""

    def test_creates_valid_xlsx(self, tmp_output_dir, sample_cui_data):
        """create_cui_xlsx produces an openable XLSX file."""
        from formatters.cui_formatter import CUIXlsxFormatter

        fmt = CUIXlsxFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['financial']
        filepath = fmt.create_cui_xlsx(doc_data, "test_financial.xlsx")

        assert os.path.exists(filepath)
        assert filepath.endswith('.xlsx')
        wb = load_workbook(filepath)
        assert len(wb.sheetnames) >= 1
        wb.close()


class TestCUIEmailFormatter:
    """Tests for CUI plain-text email generation."""

    def test_creates_valid_eml(self, tmp_output_dir, sample_cui_data):
        """create_cui_email produces a parseable EML file."""
        from formatters.cui_formatter import CUIEmailFormatter

        fmt = CUIEmailFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['critical_infrastructure']
        filepath = fmt.create_cui_email(doc_data, "test_ci.eml")

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())
        assert msg['Subject'] is not None
        assert msg['From'] is not None


class TestCUIPPTXFormatter:
    """Tests for CUI PowerPoint generation."""

    def test_creates_valid_pptx(self, tmp_output_dir, sample_cui_data):
        """create_cui_presentation produces an openable PPTX file."""
        from formatters.cui_pptx_formatter import CUIPPTXFormatter

        fmt = CUIPPTXFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['procurement']
        filepath = fmt.create_cui_presentation(doc_data, "test_procurement.pptx")

        assert os.path.exists(filepath)
        assert filepath.endswith('.pptx')
        prs = Presentation(filepath)
        assert len(prs.slides) >= 1

    def test_with_component_config(self, tmp_output_dir, sample_cui_data):
        """create_cui_presentation applies component configuration."""
        from formatters.cui_pptx_formatter import CUIPPTXFormatter
        from templates.components import ComponentMixer

        mixer = ComponentMixer(seed=42)
        config = mixer.get_random_configuration()

        fmt = CUIPPTXFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['legal']
        filepath = fmt.create_cui_presentation(doc_data, "test_styled.pptx", component_config=config)

        assert os.path.exists(filepath)
        prs = Presentation(filepath)
        assert len(prs.slides) >= 1


class TestCUIHTMLEmailFormatter:
    """Tests for CUI HTML-styled email generation."""

    def test_creates_valid_html_eml(self, tmp_output_dir, sample_cui_data):
        """create_cui_html_email produces a parseable EML with HTML part."""
        from formatters.cui_html_email_formatter import CUIHTMLEmailFormatter

        fmt = CUIHTMLEmailFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['financial']
        filepath = fmt.create_cui_html_email(doc_data, "test_html.eml")

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())

        # Should have HTML content
        parts = list(msg.walk())
        content_types = [p.get_content_type() for p in parts]
        assert 'text/html' in content_types


class TestSnykEmailGenerator:
    """Tests for Snyk vulnerability alert email generation."""

    def test_creates_valid_snyk_eml(self, tmp_output_dir):
        """create_snyk_vulnerability_alert produces a valid EML."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        gen = SnykEmailGenerator(output_dir=tmp_output_dir)
        filepath = gen.create_snyk_vulnerability_alert(
            "test_snyk.eml", is_positive=True
        )

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())
        assert msg['Subject'] is not None

    def test_negative_snyk_eml(self, tmp_output_dir):
        """Negative Snyk alert generates without CUI content."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        gen = SnykEmailGenerator(output_dir=tmp_output_dir)
        filepath = gen.create_snyk_vulnerability_alert(
            "test_snyk_neg.eml", is_positive=False
        )

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())
        assert msg['Subject'] is not None

    def test_weekly_report_positive_valid_eml(self, tmp_output_dir):
        """create_snyk_weekly_report produces a valid EML with ZTMF scoring table."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        gen = SnykEmailGenerator(output_dir=tmp_output_dir)
        filepath = gen.create_snyk_weekly_report("test_weekly_pos.eml", is_positive=True)

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())
        assert msg['Subject'] is not None
        assert 'weekly report' in msg['Subject']

        # Plain text body must contain ZTMF scoring section
        plain_body = msg.get_payload(0).get_payload(decode=True).decode()
        assert 'ZTMF APPLICATION SCORING' in plain_body
        assert 'Business Owner' in plain_body
        assert 'Pass' in plain_body and '%' in plain_body

    def test_weekly_report_negative_valid_eml(self, tmp_output_dir):
        """Negative weekly report generates without CMS org details."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        gen = SnykEmailGenerator(output_dir=tmp_output_dir)
        filepath = gen.create_snyk_weekly_report("test_weekly_neg.eml", is_positive=False)

        assert os.path.exists(filepath)
        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())
        assert msg['Subject'] is not None

    def test_weekly_report_percentages_sum_to_100(self):
        """_ztmf_app_scores always returns rows where pass+marginal+fail == 100."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        gen = SnykEmailGenerator.__new__(SnykEmailGenerator)
        gen.ZTMF_APP_NAMES = SnykEmailGenerator.ZTMF_APP_NAMES
        gen._FIRST_NAMES = SnykEmailGenerator._FIRST_NAMES
        gen._LAST_NAMES = SnykEmailGenerator._LAST_NAMES

        for _ in range(30):
            rows = gen._ztmf_app_scores(random.randint(2, 5))
            for row in rows:
                total = row['pass_pct'] + row['marginal_pct'] + row['fail_pct']
                assert total == 100, (
                    f"Percentages for '{row['app']}' sum to {total}, expected 100"
                )

    def test_weekly_report_app_names_not_yubikey(self):
        """ZTMF_APP_NAMES pool contains no reference to Yubikey Manager."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        for name in SnykEmailGenerator.ZTMF_APP_NAMES:
            assert 'Yubikey' not in name
            assert 'yubikey' not in name.lower()

    def test_weekly_report_has_both_plain_and_html(self, tmp_output_dir):
        """Weekly report email contains both plain text and HTML parts."""
        from formatters.snyk_email_generator import SnykEmailGenerator

        gen = SnykEmailGenerator(output_dir=tmp_output_dir)
        filepath = gen.create_snyk_weekly_report("test_weekly_parts.eml", is_positive=True)

        with open(filepath, 'rb') as f:
            msg = email.message_from_bytes(f.read())
        content_types = [p.get_content_type() for p in msg.get_payload()]
        assert 'text/plain' in content_types
        assert 'text/html' in content_types

        # HTML must contain the ZTMF table
        html_body = next(
            p.get_payload(decode=True).decode()
            for p in msg.get_payload()
            if p.get_content_type() == 'text/html'
        )
        assert 'ZTMF Application Scoring' in html_body
        assert 'Business Owner' in html_body


class TestPDFFormGenerators:
    """Tests for PDF form faker-data generators (EFT Authorization, FOIA MedicareAuth)."""

    def test_eft_authorization_data_has_required_fields(self):
        """generate_eft_authorization_data returns all expected form field keys."""
        import sys
        sys.path.insert(0, 'src')
        from formatters.pdf_form_populator import PDFFormPopulator

        populator = PDFFormPopulator()
        data = populator.generate_eft_authorization_data()

        # Part 1: Account Holder
        assert 'txtPayee' in data and data['txtPayee']
        assert 'txtAHStreet' in data
        assert 'txtAHCity' in data
        assert 'txtTIN' in data
        assert len(data['txtTIN']) == 9

        # Part 2: Financial Institution
        assert 'txtBankName' in data and data['txtBankName']
        assert 'txtRoutingNum' in data
        assert 'txtDepositNum' in data
        assert data['txtTypeofAccount'] in ('Checking Account', 'Savings Account')

        # Part 3: Administrative
        assert data['CMS Employee'] in ('Yes', 'No')
        assert 'txtSignature' in data

    def test_eft_authorization_data_varies_each_call(self):
        """generate_eft_authorization_data produces different data on successive calls."""
        import sys
        sys.path.insert(0, 'src')
        from formatters.pdf_form_populator import PDFFormPopulator

        populator = PDFFormPopulator()
        calls = [populator.generate_eft_authorization_data() for _ in range(10)]
        payees = {d['txtPayee'] for d in calls}
        assert len(payees) > 1, "Company names should vary across calls"

    def test_foia_medicare_auth_data_has_required_fields(self):
        """generate_foia_medicare_auth_data returns all expected form field keys."""
        import sys
        sys.path.insert(0, 'src')
        from formatters.pdf_form_populator import PDFFormPopulator

        populator = PDFFormPopulator()
        data = populator.generate_foia_medicare_auth_data()

        assert 'FirstName' in data and data['FirstName']
        assert 'LastName' in data and data['LastName']
        assert 'MedicareID' in data
        assert 'Birthdate' in data
        assert 'StreetAddress' in data
        assert 'TimeframeStart' in data
        assert 'TimeframeEnd' in data
        assert isinstance(data['ReleaseRecords'], bool)
        assert 'Signature1' in data

    def test_foia_medicare_auth_mbi_format(self):
        """MedicareID follows MBI format: starts with allowed letter, 10 chars total."""
        import sys
        sys.path.insert(0, 'src')
        from formatters.pdf_form_populator import PDFFormPopulator

        populator = PDFFormPopulator()
        for _ in range(20):
            data = populator.generate_foia_medicare_auth_data()
            mbi = data['MedicareID']
            assert len(mbi) == 10
            assert mbi[0] in 'ACDEFGHJKMNPQRTUVWXY'

    def test_eft_mapping_uses_blank_template_for_both_polarities(self):
        """EFT Authorization Form mapping uses faker population for positive and negative."""
        import sys
        sys.path.insert(0, 'src')
        from formatters.pdf_form_populator import CustomerTemplateManager

        mgr = CustomerTemplateManager()
        info = mgr.template_mappings['EFT Authorization Form']

        # Must use 'template' (single form, faker-populated) not 'template_positive' (copy mode)
        assert 'template' in info
        assert 'template_positive' not in info
        assert 'generator' in info
        assert 'saved_templates' in info['template']

    def test_foia_medicare_auth_mapping_uses_saved_templates(self):
        """FOIAMedicareAuth mapping references the saved_templates directory."""
        import sys
        sys.path.insert(0, 'src')
        from formatters.pdf_form_populator import CustomerTemplateManager

        mgr = CustomerTemplateManager()
        info = mgr.template_mappings['FOIAMedicareAuth']
        assert 'saved_templates' in info['template']
        assert 'generator' in info

    def test_llm_enrichable_set_includes_pdf_templates(self):
        """PDF templates EFT Authorization Form and FOIAMedicareAuth have LLM prompts."""
        import sys
        sys.path.insert(0, 'src')
        from generators.llm_generator import ClaudeGenerator

        gen = ClaudeGenerator.__new__(ClaudeGenerator)
        # Manually invoke prompt lookup logic via generate_template_narrative
        # (requires client — just verify prompt keys are registered)
        import inspect
        src = inspect.getsource(gen.__class__.generate_template_narrative)
        assert "'EFT Authorization Form'" in src
        assert "'FOIAMedicareAuth'" in src


class TestNegativeDocumentStructure:
    """Tests for CUI negative document structure."""

    def test_negative_docx_has_structure(self, tmp_output_dir):
        """Negative DOCX documents have expected structure."""
        from formatters.cui_formatter import CUIDocxFormatter
        from generators.cui import CUIGeneratorFactory

        gen = CUIGeneratorFactory.get_generator('legal', seed=42)
        neg_data = gen.generate_negative()

        fmt = CUIDocxFormatter(output_dir=tmp_output_dir)
        filepath = fmt.create_cui_document(neg_data, "test_neg.docx")

        assert os.path.exists(filepath)
        doc = Document(filepath)
        assert len(doc.paragraphs) > 0

    def test_classification_header_present(self, tmp_output_dir, sample_cui_data):
        """When has_cui=True and classification provided, it appears in DOCX."""
        from formatters.cui_formatter import CUIDocxFormatter

        fmt = CUIDocxFormatter(output_dir=tmp_output_dir)
        doc_data = sample_cui_data['tax']
        doc_data['has_cui'] = True
        doc_data['classification'] = 'CUI//SP-TAX'

        filepath = fmt.create_cui_document(doc_data, "test_classified.docx")
        doc = Document(filepath)
        full_text = '\n'.join(p.text for p in doc.paragraphs)
        assert 'CUI' in full_text
